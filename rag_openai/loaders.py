from __future__ import annotations
from pathlib import Path
from typing import Dict, Any, List
from .utils import read_text_file

# PDF
from pypdf import PdfReader

# OCR / Images (optionnels, mais utiles)
try:
    from pdf2image import convert_from_path  # type: ignore
except Exception:  # pragma: no cover
    convert_from_path = None

try:
    import pytesseract  # type: ignore
except Exception:  # pragma: no cover
    pytesseract = None

try:
    from PIL import Image  # type: ignore
except Exception:  # pragma: no cover
    Image = None


def load_pdf(path: Path) -> List[Dict[str, Any]]:
    out: List[Dict[str, Any]] = []
    reader = PdfReader(str(path))
    for i, page in enumerate(reader.pages):
        txt = ""
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        if not txt.strip():
            # fallback OCR si dispo
            if convert_from_path is not None and pytesseract is not None:
                try:
                    imgs = convert_from_path(str(path), first_page=i + 1, last_page=i + 1, dpi=200)
                    if imgs:
                        txt = pytesseract.image_to_string(imgs[0])
                except Exception:
                    pass
        if txt.strip():
            out.append({
                "text": txt,
                "metadata": {"source": str(path), "type": "pdf", "page": i + 1},
            })
    return out


def load_image(path: Path) -> List[Dict[str, Any]]:
    docs: List[Dict[str, Any]] = []
    caption = caption_image_with_openai(path)
    if caption:
        docs.append({
            "text": caption,
            "metadata": {"source": str(path), "type": "image", "method": "vision"},
        })
        return docs
    if Image is None or pytesseract is None:
        return []
    try:
        img = Image.open(path)
        txt = pytesseract.image_to_string(img)
        if txt.strip():
            return [{"text": txt, "metadata": {"source": str(path), "type": "image"}}]
    except Exception:
        return []
    return []


def load_docx(path: Path) -> List[Dict[str, Any]]:
    from docx import Document  # lazy import

    doc = Document(str(path))
    txt = "\n".join(p.text for p in doc.paragraphs)
    return [{"text": txt, "metadata": {"source": str(path), "type": "docx"}}]


def load_pptx(path: Path) -> List[Dict[str, Any]]:
    from pptx import Presentation  # lazy import

    prs = Presentation(str(path))
    texts: List[Dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        slide_txt: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_txt.append(shape.text)
        if slide_txt:
            texts.append({
                "text": "\n".join(slide_txt),
                "metadata": {"source": str(path), "type": "pptx", "slide": i + 1},
            })
    return texts


def load_csv(path: Path) -> List[Dict[str, Any]]:
    import pandas as pd  # lazy import

    df = None
    try:
        df = pd.read_csv(path, nrows=200)
    except Exception:
        try:
            df = pd.read_csv(path, nrows=200, sep=";")
        except Exception:
            return []
    schema = ", ".join([f"{c}({str(df[c].dtype)})" for c in df.columns])
    preview = df.head(20).to_csv(index=False)
    text = f"CSV Schema: {schema}\nPreview(20 rows):\n{preview}"
    return [{"text": text, "metadata": {"source": str(path), "type": "csv"}}]


def load_md(path: Path) -> List[Dict[str, Any]]:
    txt = read_text_file(path)
    return [{"text": txt, "metadata": {"source": str(path), "type": "md"}}]


def load_txt(path: Path) -> List[Dict[str, Any]]:
    txt = read_text_file(path)
    return [{"text": txt, "metadata": {"source": str(path), "type": "txt"}}]


def load_zip(path: Path, *, verbose: bool = False) -> List[Dict[str, Any]]:
    import zipfile, tempfile
    items: List[Dict[str, Any]] = []
    with zipfile.ZipFile(path, "r") as z:
        with tempfile.TemporaryDirectory() as tmp:
            z.extractall(tmp)
            from .loaders import ingest_directory  # local import to avoid cycle

            items += ingest_directory(Path(tmp), verbose=verbose)
    return items


def load_eml(path: Path) -> List[Dict[str, Any]]:
    # Basic .eml parser: extract headers and textual body
    from email import policy
    from email.parser import BytesParser
    import re as _re

    try:
        with open(path, "rb") as f:
            msg = BytesParser(policy=policy.default).parse(f)
    except Exception:
        return []

    headers = {
        "subject": msg.get("subject", ""),
        "from": msg.get("from", ""),
        "to": msg.get("to", ""),
        "date": msg.get("date", ""),
        "message_id": msg.get("message-id", ""),
    }

    # Prefer text/plain; fallback to stripped text/html
    body_txt: str = ""
    try:
        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type()
                disp = part.get_content_disposition()
                if disp == "attachment":
                    continue
                if ctype == "text/plain":
                    body_txt = part.get_content().strip()
                    if body_txt:
                        break
            if not body_txt:
                for part in msg.walk():
                    ctype = part.get_content_type()
                    disp = part.get_content_disposition()
                    if disp == "attachment":
                        continue
                    if ctype == "text/html":
                        html = part.get_content()
                        # crude strip of HTML tags
                        body_txt = _re.sub(r"<[^>]+>", " ", html)
                        body_txt = _re.sub(r"\s+", " ", body_txt).strip()
                        if body_txt:
                            break
        else:
            ctype = msg.get_content_type()
            if ctype == "text/plain":
                body_txt = msg.get_content().strip()
            elif ctype == "text/html":
                html = msg.get_content()
                body_txt = _re.sub(r"<[^>]+>", " ", html)
                body_txt = _re.sub(r"\s+", " ", body_txt).strip()
    except Exception:
        body_txt = ""

    header_str = " | ".join([f"{k}: {v}" for k, v in headers.items() if v])
    full_text = f"{header_str}\n\n{body_txt}".strip()
    if not full_text:
        return []
    return [{
        "text": full_text,
        "metadata": {"source": str(path), "type": "eml", **headers},
    }]


LOADER_REGISTRY = {
    ".pdf": load_pdf,
    ".png": load_image,
    ".jpg": load_image,
    ".jpeg": load_image,
    ".webp": load_image,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".csv": load_csv,
    ".md": load_md,
    ".txt": load_txt,
    ".zip": load_zip,
    ".eml": load_eml,
}


def ingest_file(path: Path, *, verbose: bool = False) -> List[Dict[str, Any]]:
    loader = LOADER_REGISTRY.get(path.suffix.lower())
    if loader is None:
        if verbose:
            print(f"[ingest:skip] {path} (unsupported)")
        return []
    try:
        items = loader(path) if loader is not load_zip else loader(path, verbose=verbose)
        if verbose:
            print(f"[ingest:file] {path} -> {len(items)} logical docs")
        return items
    except Exception as e:
        if verbose:
            print(f"[ingest:error] {path} -> {e}")
        return []


def ingest_directory(root: Path, *, verbose: bool = False) -> List[Dict[str, Any]]:
    docs: List[Dict[str, Any]] = []
    if verbose:
        print(f"[ingest:start] Scanning {root}")
    for p in root.rglob("*"):
        if p.is_file():
            docs += ingest_file(p, verbose=verbose)
    if verbose:
        print(f"[ingest:done] Total logical documents: {len(docs)}")
    return docs

    import base64
from openai import OpenAI
from .config import OPENAI_API_KEY, OPENAI_VISION_MODEL

def caption_image_with_openai(path: Path) -> str:
    if not OPENAI_API_KEY:
        return ""
    client = OpenAI(api_key=OPENAI_API_KEY)
    with open(path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode("utf-8")
    uri = f"data:image/{path.suffix.lstrip('.').lower()};base64,{b64}"
    messages = [
        {"role": "system", "content": "You are a helpful assistant that describes images."},
        {
            "role": "user",
            "content": [
                {"type": "image_url", "image_url": {"url": uri}},
                {"type": "text", "text": "Provide a concise description of this image in one or two sentences."},
            ],
        },
    ]
    resp = client.chat.completions.create(
        model=OPENAI_VISION_MODEL,
        messages=messages,
        temperature=0.0,
        max_tokens=150,
    )
    return resp.choices[0].message.content.strip()
